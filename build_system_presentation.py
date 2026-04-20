from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.units import cm
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parent
PPTX_PATH = ROOT / "PTR_System_Presentation.pptx"
PDF_PATH = ROOT / "PTR_System_Presentation.pdf"
PGP_LOGO = ROOT / "PGP.png"
PHO_LOGO = ROOT / "PHO.png"


SLIDES = [
    (
        "Property Stock Transfer Report System",
        [
            "Provincial Health Office - Palawan",
            "Web-based inventory and PTR workflow",
            f"Generated: {datetime.now().strftime('%Y-%m-%d')}",
        ],
    ),
    (
        "System Overview",
        [
            "Built with PHP, MySQL, Bootstrap, and JavaScript",
            "Supports role-based access (Admin and Encoder)",
            "Tracks stock movement and PTR lifecycle",
            "Includes reports, print previews, and transaction history",
        ],
    ),
    (
        "Core Modules",
        [
            "Dashboard and inventory summary",
            "Create PTR and pending transaction release flow",
            "Transaction History with grouped PTR view",
            "Current Stock Report, Stock Card, and Incident Report",
            "Notifications and master item management",
        ],
    ),
    (
        "Create PTR Workflow",
        [
            "Users add multiple line items with batch, quantity, and unit cost",
            "System computes totals and prepares print preview",
            "Drafts can be saved as pending before release",
            "Release updates stock and records PTR entries",
        ],
    ),
    (
        "Transaction History & Print",
        [
            "Grouped display by PTR number for clean review",
            "Edit signatory names in a compact panel",
            "Print button opens formatted A4 landscape output",
            "Signatories can be saved and reused per PTR number",
        ],
    ),
    (
        "Signatory Name Persistence",
        [
            "Prepared by, Approved by, and Issued by are now editable and savable",
            "Saved values are stored in ptr_print_signatories table",
            "Defaults load automatically when no saved value exists",
            "Create PTR and Transaction History share the same saved values",
        ],
    ),
    (
        "Reports & Analytics",
        [
            "Current stock by item and batch",
            "Stock card movement history",
            "Outbound summary for released quantities",
            "Incident report logging and review",
        ],
    ),
    (
        "Security and Access Control",
        [
            "Session-based login and route protection",
            "Role checks for module access restrictions",
            "Encoder mutation guard blocks restricted actions",
            "Input validation and server-side checks on save/update flows",
        ],
    ),
    (
        "Benefits to Operations",
        [
            "Faster PTR preparation and consistent print outputs",
            "Reduced manual encoding errors",
            "Better traceability for stock transfers",
            "Centralized records for audit and monitoring",
        ],
    ),
    (
        "Next Improvements",
        [
            "Export-ready dashboards and KPIs",
            "Automated email notifications for approvals/releases",
            "Advanced audit trail and user activity logs",
            "Mobile-friendly optimization for field usage",
        ],
    ),
]


def add_logos(slide):
    if PGP_LOGO.exists():
        slide.shapes.add_picture(str(PGP_LOGO), Inches(0.3), Inches(0.2), height=Inches(0.6))
    if PHO_LOGO.exists():
        slide.shapes.add_picture(str(PHO_LOGO), Inches(12.2), Inches(0.2), height=Inches(0.6))


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[1]

    for title, bullets in SLIDES:
        slide = prs.slides.add_slide(layout)
        add_logos(slide)
        slide.shapes.title.text = title
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.size = Pt(34)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(15, 64, 48)

        body = slide.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(22 if i == 0 and len(bullets) <= 3 else 20)
            p.font.color.rgb = RGBColor(36, 36, 36)
    prs.save(PPTX_PATH)


def build_pdf():
    c = canvas.Canvas(str(PDF_PATH), pagesize=landscape(A4))
    width, height = landscape(A4)
    for idx, (title, bullets) in enumerate(SLIDES, start=1):
        c.setFont("Helvetica-Bold", 24)
        c.setFillColorRGB(0.06, 0.28, 0.22)
        c.drawString(2 * cm, height - 2.2 * cm, title)

        y = height - 4 * cm
        c.setFillColorRGB(0.12, 0.12, 0.12)
        for b in bullets:
            c.setFont("Helvetica", 15)
            c.drawString(2.5 * cm, y, f"- {b}")
            y -= 1.2 * cm

        c.setFont("Helvetica-Oblique", 10)
        c.drawRightString(width - 1.5 * cm, 1.0 * cm, f"Slide {idx} of {len(SLIDES)}")
        c.showPage()
    c.save()


if __name__ == "__main__":
    build_pptx()
    build_pdf()
    print(f"Created: {PPTX_PATH}")
    print(f"Created: {PDF_PATH}")
